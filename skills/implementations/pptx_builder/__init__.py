"""
PPTX Builder Skill
Create and modify PowerPoint presentations using python-pptx.
"""
from __future__ import annotations

from pathlib import Path
from typing import Any, Optional

from nexus.skills.registry import BaseSkill, SkillMeta

SKILL_META = SkillMeta(
    name="pptx_builder",
    description="Create and edit PowerPoint presentations: "
                "add slides, insert text/images/tables, apply themes, save as .pptx",
    version="1.0.0",
    domains=["creative", "operations", "research"],
    triggers=["powerpoint", "pptx", "presentation", "slide", "投影片", "簡報"],
    requires=["python-pptx"],
    is_local=True,
)


class Skill(BaseSkill):
    meta = SKILL_META

    async def run(
        self,
        operation: str,          # create | add_slide | add_text | add_image | add_table | save
        file_path: str = "",
        output_path: str = "",
        **kwargs,
    ) -> Any:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            return {"error": "python-pptx not installed. Run: pip install python-pptx"}

        op = operation.lower()

        # ── Create new presentation ───────────────────────────────────────────
        if op == "create":
            prs = Presentation()
            # Widescreen 16:9
            prs.slide_width  = Inches(13.33)
            prs.slide_height = Inches(7.5)
            save_to = output_path or file_path or "presentation.pptx"
            Path(save_to).parent.mkdir(parents=True, exist_ok=True)
            prs.save(save_to)
            return {"created": save_to, "slide_count": 0}

        # ── Load existing presentation ────────────────────────────────────────
        def _load(path: str) -> "Presentation":
            if path and Path(path).exists():
                return Presentation(path)
            prs = Presentation()
            prs.slide_width  = Inches(13.33)
            prs.slide_height = Inches(7.5)
            return prs

        # ── Add slide ─────────────────────────────────────────────────────────
        if op == "add_slide":
            prs = _load(file_path)
            layout_idx = kwargs.get("layout", 1)  # 1 = title + content
            layout = prs.slide_layouts[min(layout_idx, len(prs.slide_layouts) - 1)]
            slide  = prs.slides.add_slide(layout)

            title   = kwargs.get("title", "")
            content = kwargs.get("content", "")
            if title and slide.shapes.title:
                slide.shapes.title.text = title
            if content and len(slide.placeholders) > 1:
                slide.placeholders[1].text = content

            save_to = output_path or file_path or "presentation.pptx"
            prs.save(save_to)
            return {
                "saved":       save_to,
                "slide_count": len(prs.slides),
                "slide_index": len(prs.slides) - 1,
            }

        # ── Add text box ──────────────────────────────────────────────────────
        elif op == "add_text":
            prs   = _load(file_path)
            idx   = kwargs.get("slide_index", len(prs.slides) - 1)
            if not prs.slides:
                return {"error": "No slides in presentation"}
            slide = prs.slides[min(idx, len(prs.slides) - 1)]

            left   = Inches(kwargs.get("left",   1.0))
            top    = Inches(kwargs.get("top",    1.0))
            width  = Inches(kwargs.get("width",  8.0))
            height = Inches(kwargs.get("height", 1.5))
            text   = kwargs.get("text", "")
            font_size = kwargs.get("font_size", 18)
            bold   = kwargs.get("bold", False)
            color  = kwargs.get("color", "000000")

            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf    = txBox.text_frame
            tf.word_wrap = True
            p  = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = RGBColor.from_string(color)

            save_to = output_path or file_path or "presentation.pptx"
            prs.save(save_to)
            return {"saved": save_to, "text_added": text[:50]}

        # ── Add image ─────────────────────────────────────────────────────────
        elif op == "add_image":
            prs = _load(file_path)
            idx = kwargs.get("slide_index", len(prs.slides) - 1)
            if not prs.slides:
                return {"error": "No slides in presentation"}
            slide      = prs.slides[min(idx, len(prs.slides) - 1)]
            image_path = kwargs.get("image_path", "")
            if not image_path or not Path(image_path).exists():
                return {"error": f"Image not found: {image_path}"}

            left   = Inches(kwargs.get("left",   1.0))
            top    = Inches(kwargs.get("top",    1.5))
            width  = Inches(kwargs.get("width",  6.0))
            slide.shapes.add_picture(image_path, left, top, width=width)

            save_to = output_path or file_path or "presentation.pptx"
            prs.save(save_to)
            return {"saved": save_to, "image": image_path}

        # ── Add table ─────────────────────────────────────────────────────────
        elif op == "add_table":
            prs = _load(file_path)
            idx = kwargs.get("slide_index", len(prs.slides) - 1)
            if not prs.slides:
                return {"error": "No slides in presentation"}
            slide = prs.slides[min(idx, len(prs.slides) - 1)]
            data  = kwargs.get("data", [])   # list of lists
            if not data:
                return {"error": "data required for table"}

            rows = len(data)
            cols = max(len(r) for r in data)
            left  = Inches(kwargs.get("left",  1.0))
            top   = Inches(kwargs.get("top",   2.0))
            width = Inches(kwargs.get("width", 10.0))
            height= Inches(kwargs.get("height", 0.5 * rows))

            table = slide.shapes.add_table(rows, cols, left, top, width, height).table
            for r_idx, row in enumerate(data):
                for c_idx, val in enumerate(row):
                    cell = table.cell(r_idx, c_idx)
                    cell.text = str(val)
                    if r_idx == 0:  # header row bold
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                run.font.bold = True

            save_to = output_path or file_path or "presentation.pptx"
            prs.save(save_to)
            return {"saved": save_to, "table_shape": (rows, cols)}

        # ── Read slide content ────────────────────────────────────────────────
        elif op == "read":
            prs = _load(file_path)
            slides_data = []
            for i, slide in enumerate(prs.slides):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        texts.append(shape.text.strip())
                slides_data.append({"index": i, "texts": texts})
            return {"slide_count": len(prs.slides), "slides": slides_data}

        return {"error": f"Unknown operation: {operation}"}
